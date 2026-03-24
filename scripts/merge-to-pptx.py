#!/usr/bin/env python3
import os
import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def parse_args():
    args = sys.argv[1:]
    dir_path = ""
    output = None

    i = 0
    while i < len(args):
        if args[i] in ["--output", "-o"]:
            output = args[i + 1]
            i += 2
        elif not args[i].startswith("-"):
            dir_path = args[i]
            i += 1
        else:
            i += 1

    if not dir_path:
        print("Usage: python merge-to-pptx.py <slide-deck-dir> [--output filename.pptx]", file=sys.stderr)
        sys.exit(1)

    return dir_path, output

def find_slide_images(dir_path):
    if not os.path.exists(dir_path):
        print(f"Directory not found: {dir_path}", file=sys.stderr)
        sys.exit(1)

    slides_dir = os.path.join(dir_path, "slides")
    search_dir = slides_dir if os.path.exists(slides_dir) else dir_path

    slide_pattern = re.compile(r'^(\d+)-slide-.*\.(png|jpg|jpeg)$', re.IGNORECASE)
    prompts_dir = os.path.join(dir_path, "prompts")
    has_prompts = os.path.exists(prompts_dir)

    slides = []
    for filename in os.listdir(search_dir):
        match = slide_pattern.match(filename)
        if match:
            index = int(match.group(1))
            base_name = re.sub(r'\.(png|jpg|jpeg)$', '', filename, flags=re.IGNORECASE)
            prompt_path = os.path.join(prompts_dir, f"{base_name}.md") if has_prompts else None

            slides.append({
                'filename': filename,
                'path': os.path.join(search_dir, filename),
                'index': index,
                'prompt_path': prompt_path if prompt_path and os.path.exists(prompt_path) else None
            })

    slides.sort(key=lambda x: x['index'])

    if not slides:
        print(f"No slide images found in: {search_dir}", file=sys.stderr)
        print("Expected format: 01-slide-*.png, 02-slide-*.png, etc.", file=sys.stderr)
        sys.exit(1)

    return slides

def find_base_prompt():
    script_dir = Path(__file__).parent
    base_prompt_path = script_dir.parent / "references" / "base-prompt.md"
    if base_prompt_path.exists():
        return base_prompt_path.read_text(encoding='utf-8')
    return None

def create_pptx(slides, output_path):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    base_prompt = find_base_prompt()
    notes_count = 0

    for slide_info in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        pic = slide.shapes.add_picture(
            slide_info['path'],
            0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        if slide_info['prompt_path']:
            with open(slide_info['prompt_path'], 'r', encoding='utf-8') as f:
                slide_prompt = f.read()
            full_notes = f"{base_prompt}\n\n---\n\n{slide_prompt}" if base_prompt else slide_prompt
            slide.notes_slide.notes_text_frame.text = full_notes
            notes_count += 1

        print(f"Added: {slide_info['filename']}{' (with notes)' if slide_info['prompt_path'] else ''}")

    prs.save(output_path)
    print(f"\nCreated: {output_path}")
    print(f"Total slides: {len(slides)}")
    if notes_count > 0:
        print(f"Slides with notes: {notes_count}{' (includes base prompt)' if base_prompt else ''}")

def main():
    dir_path, output = parse_args()
    slides = find_slide_images(dir_path)

    dir_name = os.path.basename(dir_path)
    if dir_name == "slide-deck":
        dir_name = os.path.basename(os.path.dirname(dir_path))

    output_path = output or os.path.join(dir_path, f"{dir_name}.pptx")

    print(f"Found {len(slides)} slides in: {dir_path}\n")
    create_pptx(slides, output_path)

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
